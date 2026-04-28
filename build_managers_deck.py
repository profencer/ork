"""Build ork-for-managers-ru.pptx — visual identity matches ork-current-state-ru.pptx.

Audience: technical / IT managers. Format: slides in Russian that map ork's
technical capabilities onto enterprise business value, using the competitor
report as the framing matrix.
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# Visual identity (lifted from ork-current-state-ru.pptx)
SLIDE_W = 12191695
SLIDE_H = 6858000
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_LIGHT = RGBColor(0xF5, 0xF1, 0xEA)
ACCENT = RGBColor(0xE6, 0x4B, 0x2A)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x55, 0x60, 0x77)
HAIRLINE = RGBColor(0xCB, 0xC2, 0xB2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

LEFT = 548640
RIGHT_EDGE = SLIDE_W - LEFT
HEADER_LABEL_Y = 457200
TITLE_Y = 713232
DIVIDER_Y = 1325880
CONTENT_Y = 1554480
FOOTER_Y = 6400800
FOOTER_H = 274320

FONT = "Calibri"
TOTAL_SLIDES = 18  # update if slide count changes
BRAND_FOOTER = "ork — мульти-агентная оркестрация на Rust"


def _set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, rgb, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    _set_fill(s, rgb)
    return s


def _add_text(slide, x, y, w, h, text, size_pt, *, bold=False, color=INK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, items, size_pt=11, color=INK,
                 line_spacing=1.25, marker="·  "):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            head, tail = item
            r1 = p.add_run()
            r1.text = marker + head
            r1.font.name = FONT
            r1.font.size = Pt(size_pt)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = " — " + tail
            r2.font.name = FONT
            r2.font.size = Pt(size_pt)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = marker + item
            r.font.name = FONT
            r.font.size = Pt(size_pt)
            r.font.color.rgb = color
    return tb


def _add_line(slide, x1, y1, x2, y2, rgb=HAIRLINE, weight_pt=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = rgb
    line.line.width = Pt(weight_pt)
    return line


def _round_card(slide, x, y, w, h, *, accent_rgb=None, accent_w=109728):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _set_fill(card, CARD)
    # tighten the corner radius a hair
    try:
        card.adjustments[0] = 0.06
    except Exception:
        pass
    if accent_rgb is not None:
        _add_rect(slide, x, y, accent_w, h, accent_rgb)
    return card


def _add_header(slide, label, title, page_idx):
    # background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_LIGHT)
    # accent square next to label
    _add_rect(slide, LEFT, 502920, 109728, 457200, ACCENT)
    # ALL CAPS small label
    _add_text(slide, 777240, HEADER_LABEL_Y, 10058400, 274320,
              label, 11, bold=True, color=ACCENT)
    # large title
    _add_text(slide, 777240, TITLE_Y, 10972800, 548640,
              title, 28, bold=True, color=INK)
    # divider line
    _add_line(slide, LEFT, DIVIDER_Y, LEFT + 11064240, DIVIDER_Y, MUTED, 0.75)
    # footer
    _add_text(slide, LEFT, FOOTER_Y, 5486400, FOOTER_H,
              BRAND_FOOTER, 9, color=MUTED)
    _add_text(slide, 10058400, FOOTER_Y, 1828800, FOOTER_H,
              f"{page_idx} / {TOTAL_SLIDES}", 9, color=MUTED, align=PP_ALIGN.RIGHT)


# -------------------------- slide builders --------------------------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)
    _add_rect(s, LEFT, 2286000, 164592, 2286000, ACCENT)
    _add_text(s, 868680, 2194560, 10058400, 457200,
              "Мульти-агентная оркестрация для enterprise",
              14, color=WHITE)
    _add_text(s, 868680, 2606040, 10058400, 1280160,
              "ork", 88, bold=True, color=WHITE)
    _add_text(s, 868680, 4114800, 10058400, 548640,
              "Платформа агентного AI: для технических руководителей",
              24, color=WHITE)
    _add_text(s, 868680, 4846320, 10058400, 365760,
              "A2A-first · MCP · Kong + Kafka · Rust 2024 · ADR-driven",
              14, color=ACCENT)
    _add_text(s, 868680, 6217920, 10058400, 365760,
              "Апрель 2026", 12, color=WHITE)


def slide_context(prs):
    _add_header_and_index(prs, "КОНТЕКСТ", "Зачем enterprise-у платформа агентного AI", 2)
    s = prs.slides[-1]
    intro = ("Агентный AI выходит из пилотов в production. У IT-руководителя "
             "появляется четыре одновременных требования к платформе: данные, "
             "безопасность, стоимость и скорость доставки.")
    _add_text(s, 777240, CONTENT_Y, 10515600, 640080, intro, 14, color=INK)

    cards = [
        ("Данные и заземление",
         "Агенты должны видеть актуальные корпоративные данные "
         "(БД, API, события), иначе hallucinations превращаются в инциденты."),
        ("Безопасность и комплаенс",
         "SSO, RBAC, аудит, изоляция тенантов — без этого AI "
         "не пускают в продуктив."),
        ("Стоимость и предсказуемость",
         "Токены и compute растут нелинейно. Нужен контроль контекста, "
         "кэширование и выбор провайдера под задачу."),
        ("Скорость pilot → production",
         "Бизнес ждёт работающие сценарии за недели, а не кварталы. "
         "Нужны и no-code, и code-first контуры."),
    ]
    _grid_2x2(s, cards, top_y=2331720)


def slide_what_is_ork(prs):
    _add_header_and_index(prs, "ОБЗОР", "Что такое ork в одном слайде", 3)
    s = prs.slides[-1]
    lead = ("ork — это open-платформа оркестрации множества AI-агентов на Rust 2024. "
            "Один общий протокол A2A, гибридный транспорт Kong + Kafka, MCP как единая "
            "шина инструментов и гексагональная архитектура с изолированной доменной логикой.")
    _add_text(s, 777240, CONTENT_Y, 10515600, 640080, lead, 14, color=INK)

    pillars = [
        ("A2A-first",
         "Все агенты — локальные и удалённые — говорят одним протоколом: "
         "cards, tasks, типизированные parts, streaming, cancel, push."),
        ("Гибридный транспорт",
         "Sync через Kong (HTTP/SSE). Async через Kafka. Нет проприетарного брокера, "
         "нет vendor lock-in на инфраструктурном уровне."),
        ("MCP как шина инструментов",
         "Внешние интеграции — стандартные MCP-серверы. Внутренние — нативные Rust-tools. "
         "Единый каталог, контролируемый scope."),
        ("Гексагональная архитектура",
         "Доменная логика изолирована от транспорта и БД. Замена Postgres, Redis, Kafka — "
         "через порты, без переписывания ядра."),
    ]
    _grid_2x2(s, pillars, top_y=2331720)


def slide_value_overview(prs):
    _add_header_and_index(prs, "БИЗНЕС-КАРТА",
                          "Восемь enterprise-требований и как ork их закрывает", 4)
    s = prs.slides[-1]
    _add_text(s, 777240, CONTENT_Y, 10515600, 365760,
              "Каждое требование — отдельный слайд далее. ADR в скобках.",
              12, color=MUTED)

    features = [
        ("Данные в реальном времени", "Generic Gateways · MCP · Kafka (ADR 0010, 0013)"),
        ("Безопасность и комплаенс", "Tenant security · RBAC scopes (ADR 0020, 0021)"),
        ("Vendor-агностичность", "Multi-LLM каталог (ADR 0012)"),
        ("Стоимость и токены ↓", "Workflow engine · Dynamic embeds (ADR 0015, 0018)"),
        ("Масштаб и устойчивость", "Kong + Kafka · параллельные DAG"),
        ("Observability и governance", "Workflow viewer · трассировка (ADR 0022)"),
        ("Каналы доступа", "Slack · Teams · REST · WebUI (ADR 0013, 0017)"),
        ("Скорость экспериментов", "YAML-workflow · WebUI-чат (ADR 0017)"),
    ]
    _grid_4x2(s, features, top_y=2057400)


def slide_feature(prs, page, label, title, kicker, tech_items, value_items):
    """One enterprise-feature slide: tech-side card vs. business-value card."""
    _add_header_and_index(prs, label, title, page)
    s = prs.slides[-1]
    if kicker:
        _add_text(s, 777240, CONTENT_Y, 10515600, 548640, kicker, 13, color=MUTED)

    card_y = 2103120
    card_h = 4023360
    gap = 182880
    card_w = (RIGHT_EDGE - LEFT - gap) // 2

    # left card: tech
    _round_card(s, LEFT + 228600, card_y, card_w, card_h, accent_rgb=INK)
    _add_text(s, LEFT + 228600 + 274320, card_y + 182880, card_w - 411480, 411480,
              "Реализация в ork", 16, bold=True, color=INK)
    _add_text(s, LEFT + 228600 + 274320, card_y + 594360, card_w - 411480, 274320,
              "технический контур", 10, color=MUTED)
    _add_bullets(s, LEFT + 228600 + 274320, card_y + 914400,
                 card_w - 411480, card_h - 1097280,
                 tech_items, size_pt=12)

    # right card: business value
    rx = LEFT + 228600 + card_w + gap
    _round_card(s, rx, card_y, card_w, card_h, accent_rgb=ACCENT)
    _add_text(s, rx + 274320, card_y + 182880, card_w - 411480, 411480,
              "Бизнес-ценность", 16, bold=True, color=INK)
    _add_text(s, rx + 274320, card_y + 594360, card_w - 411480, 274320,
              "что получает enterprise IT", 10, color=MUTED)
    _add_bullets(s, rx + 274320, card_y + 914400,
                 card_w - 411480, card_h - 1097280,
                 value_items, size_pt=12)


def slide_differentiators(prs):
    _add_header_and_index(prs, "ОТЛИЧИЯ", "Чем ork отличается от альтернатив", 13)
    s = prs.slides[-1]
    _add_text(s, 777240, CONTENT_Y, 10515600, 548640,
              "Сравнение позиционирования, не feature-by-feature таблица.",
              13, color=MUTED)

    diffs = [
        ("Rust 2024 в ядре",
         "Память безопасна по построению, без GC-пауз. Низкий overhead "
         "на агента — важно при сотнях параллельных задач."),
        ("Open-source ядро",
         "Нет проприетарных протоколов в основе. A2A — открытый стандарт, "
         "MCP — открытый стандарт. Платформа не блокирует данные."),
        ("ADR-driven эволюция",
         "Каждое архитектурное решение зафиксировано как ADR с критериями приёмки. "
         "27 ADR на момент апреля 2026 — прозрачная история выбора."),
        ("Гексагональная изоляция",
         "Postgres, Redis, Kafka, LLM-провайдер — заменяемые адаптеры. "
         "Доменное ядро не переписывается при смене инфраструктуры."),
    ]
    _grid_2x2(s, diffs, top_y=2331720)


def slide_maturity(prs):
    _add_header_and_index(prs, "ЗРЕЛОСТЬ", "Что уже работает в коде сегодня", 14)
    s = prs.slides[-1]

    items = [
        ("Полный A2A-стек",
         "cards · tasks · streaming · cancel · push (ADR 0003, 0008, 0009)"),
        ("Гибридный транспорт",
         "Kong/HTTP+SSE для sync, Kafka для async (ADR 0004)"),
        ("Workflow engine",
         "DAG-исполнитель, embed-резолв, артефакты (ADR 0015, 0016, 0018)"),
        ("Tool plane",
         "native + LLM-tools + MCP (rmcp 0.16, ADR 0010, 0011)"),
        ("Multi-LLM каталог",
         "OpenAI-совместимый формат, провайдеры заменяемы (ADR 0012)"),
        ("Generic Gateways",
         "REST · webhook · WebUI · MCP-as-gateway (ADR 0013)"),
        ("Push-уведомления",
         "ES256 JWS + envelope-шифрование (ADR 0009)"),
        ("Артефакты",
         "S3 / MinIO с версионируемым append-API (ADR 0016)"),
        ("Federation",
         "A2aRemoteAgent + Python-пир в demo (LangGraph, ADR 0007)"),
        ("Web UI",
         "React-чат-клиент поверх A2A-стримов (ADR 0017)"),
    ]
    # two columns of 5
    col_w = 5349240
    col_x_left = 777240
    col_x_right = 6400800
    top = CONTENT_Y
    _add_bullets(s, col_x_left, top, col_w, 4500000,
                 items[:5], size_pt=12, line_spacing=1.4)
    _add_bullets(s, col_x_right, top, col_w, 4500000,
                 items[5:], size_pt=12, line_spacing=1.4)
    _add_text(s, LEFT, 6126480, 10515600, 274320,
              "Демо в demo/ — 11 стадий, ~15 минут, всё на localhost.",
              11, color=MUTED)


def slide_roadmap(prs):
    _add_header_and_index(prs, "ROADMAP", "Что в очереди ADR", 15)
    s = prs.slides[-1]
    _add_text(s, 777240, CONTENT_Y, 10515600, 548640,
              "Приоритеты на следующие циклы. Каждый пункт — отдельный ADR с критериями приёмки.",
              13, color=MUTED)

    items = [
        ("ADR 0019 · Scheduled tasks",
         "Cron-планировщик для агентов: периодические workflows и SLA-окна."),
        ("ADR 0020 · Tenant security & trust",
         "Изоляция тенантов на уровне данных, секретов и сетевого периметра."),
        ("ADR 0021 · RBAC scopes",
         "Granular scopes для пользователей, агентов и API-токенов."),
        ("ADR 0022 · Observability",
         "OpenTelemetry-трассировка, метрики токенов и стоимости, аудит."),
        ("ADR 0023 · Migration & rollout",
         "Контракт обратной совместимости и стратегия выкатки в production."),
        ("ADR 0024–0027 · Расширения",
         "WASM-плагины, typed output validation, topology selection, human-in-the-loop."),
    ]
    _grid_2x3(s, items, top_y=2240280)


def slide_business_summary(prs):
    _add_header_and_index(prs, "БИЗНЕС-ЭФФЕКТ", "Сводный эффект для IT-бюджета", 16)
    s = prs.slides[-1]

    metrics = [
        ("TCO платформы ↓",
         "Rust + filtered context + multi-LLM маршрутизация — снижение затрат "
         "на токены и compute на каждый сценарий."),
        ("Time-to-pilot недели, не кварталы",
         "YAML-workflow + WebUI-чат + готовые gateway-адаптеры сокращают цикл "
         "от идеи до работающего пилота."),
        ("Vendor lock-in → 0",
         "Open A2A, open MCP, заменяемые адаптеры БД и LLM. Инвестиции в логику "
         "переживают смену поставщика."),
        ("Audit и governance ✓",
         "Workflow viewer, трассировка вызовов, RBAC-scopes — единая точка "
         "контроля для compliance и SRE."),
    ]
    _grid_2x2(s, metrics, top_y=2057400)


def slide_demo(prs):
    _add_header_and_index(prs, "ДЕМО", "15 минут на localhost: что показать руководителю", 17)
    s = prs.slides[-1]
    _add_text(s, 777240, CONTENT_Y, 10515600, 548640,
              "Три сцены, которые отвечают на вопросы менеджера: «реально ли работает», "
              "«как масштабируется», «как интегрируется».",
              13, color=MUTED)

    scenes = [
        ("Сцена 1 · A2A в действии",
         "Агент-карточки, типизированные parts, SSE-стриминг ответа. "
         "Доказательство: единый протокол вместо зоопарка SDK."),
        ("Сцена 2 · Federation + LangGraph",
         "Rust-оркестратор зовёт Python-агента на LangGraph через A2A. "
         "Доказательство: интеграция со стеком данных-сайентистов без перевода."),
        ("Сцена 3 · MCP + push + артефакты",
         "Workflow дергает внешний MCP-tool, получает артефакт, шлёт подписанный push. "
         "Доказательство: end-to-end production-сценарий."),
    ]
    _grid_3_rows(s, scenes, top_y=2057400)


def slide_call_to_action(prs):
    _add_header_and_index(prs, "ЧТО ДАЛЬШЕ", "Как двигаться дальше", 18)
    s = prs.slides[-1]

    steps = [
        ("Шаг 1 · 30 минут",
         "Прогон демо на localhost. Проверка A2A, federation, MCP, push."),
        ("Шаг 2 · 1 неделя",
         "Pilot-сценарий из реального бэклога: один workflow, один MCP-tool, один gateway."),
        ("Шаг 3 · 1 квартал",
         "Production-выкатка с RBAC, observability и SLA: ADR 0020–0023 в работе."),
    ]
    _grid_3_rows(s, steps, top_y=CONTENT_Y + 274320)

    _add_text(s, 777240, 5760720, 10515600, 365760,
              "ork — open code, ADR-driven, готов к интеграции с вашим IT-ландшафтом.",
              14, bold=True, color=ACCENT)


# -------------------------- grid helpers --------------------------

def _add_header_and_index(prs, label, title, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(s, label, title, page)


def _grid_2x2(slide, items, top_y):
    card_w = 5349240
    card_h = 1783080
    gap_x = 274320
    gap_y = 228600
    xs = [777240, 777240 + card_w + gap_x]
    ys = [top_y, top_y + card_h + gap_y]
    for i, (head, body) in enumerate(items[:4]):
        x = xs[i % 2]
        y = ys[i // 2]
        _round_card(slide, x, y, card_w, card_h)
        _add_text(slide, x + 274320, y + 182880, card_w - 411480, 411480,
                  head, 16, bold=True, color=INK)
        _add_text(slide, x + 274320, y + 640080, card_w - 411480, card_h - 822960,
                  body, 12, color=INK)


def _grid_4x2(slide, items, top_y):
    card_w = 2606040
    card_h = 1554480
    gap_x = 137160
    gap_y = 137160
    base_x = 777240
    for i, (head, body) in enumerate(items[:8]):
        col = i % 4
        row = i // 4
        x = base_x + col * (card_w + gap_x)
        y = top_y + row * (card_h + gap_y)
        _round_card(slide, x, y, card_w, card_h, accent_rgb=ACCENT, accent_w=73152)
        _add_text(slide, x + 182880, y + 182880, card_w - 320040, 411480,
                  head, 13, bold=True, color=INK)
        _add_text(slide, x + 182880, y + 594360, card_w - 320040, card_h - 731520,
                  body, 10, color=MUTED)


def _grid_2x3(slide, items, top_y):
    card_w = 5349240
    card_h = 1188720
    gap_x = 274320
    gap_y = 137160
    xs = [777240, 777240 + card_w + gap_x]
    for i, (head, body) in enumerate(items[:6]):
        col = i % 2
        row = i // 2
        x = xs[col]
        y = top_y + row * (card_h + gap_y)
        _round_card(slide, x, y, card_w, card_h)
        _add_text(slide, x + 274320, y + 137160, card_w - 411480, 365760,
                  head, 13, bold=True, color=INK)
        _add_text(slide, x + 274320, y + 502920, card_w - 411480, card_h - 640080,
                  body, 11, color=INK)


def _grid_3_rows(slide, items, top_y):
    card_w = 10972800
    card_h = 1097280
    gap_y = 137160
    for i, (head, body) in enumerate(items[:3]):
        x = 777240
        y = top_y + i * (card_h + gap_y)
        _round_card(slide, x, y, card_w, card_h, accent_rgb=ACCENT)
        _add_text(slide, x + 320040, y + 137160, card_w - 457200, 365760,
                  head, 14, bold=True, color=INK)
        _add_text(slide, x + 320040, y + 502920, card_w - 457200, card_h - 640080,
                  body, 12, color=INK)


# -------------------------- main --------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    slide_cover(prs)
    # 2. Context
    slide_context(prs)
    # 3. What is ork
    slide_what_is_ork(prs)
    # 4. Value overview (8 features)
    slide_value_overview(prs)

    # 5. Real-time data integration
    slide_feature(
        prs, 5, "ТРЕБОВАНИЕ 1 · ДАННЫЕ В РЕАЛЬНОМ ВРЕМЕНИ",
        "Заземление агентов на корпоративные данные",
        "Hallucinations возникают там, где агент работает на устаревшем или урезанном "
        "контексте. ork даёт три канала к живым данным.",
        [
            ("Generic Gateways (ADR 0013)",
             "REST, webhook, MCP-as-gateway, WebUI — единый адаптерный слой"),
            ("MCP-инструменты (ADR 0010)",
             "БД, API, файловые системы как стандартные MCP-серверы"),
            ("Kafka event mesh",
             "async-канал для событий и reactive-агентов"),
            ("Postgres + Redis",
             "состояние и кэш через порты ork-persistence / ork-cache"),
        ],
        [
            ("Точность ответов ↑",
             "агент видит актуальные данные, а не снимок месячной давности"),
            ("Снижение инцидентов",
             "меньше hallucinations — меньше ложных автоматических действий"),
            ("Скорость интеграции",
             "новый источник = новый MCP-сервер, без правки ядра"),
            ("Работа с legacy",
             "REST-gateway оборачивает любую внутреннюю систему"),
        ],
    )

    # 6. Security & compliance
    slide_feature(
        prs, 6, "ТРЕБОВАНИЕ 2 · БЕЗОПАСНОСТЬ И КОМПЛАЕНС",
        "Корпоративная безопасность как первый класс",
        "Безопасность не дорабатывается потом — она спроектирована в ADR "
        "0020 (tenant) и 0021 (RBAC) и в самом A2A-протоколе.",
        [
            ("SSO / OAuth2",
             "интеграция с Azure AD, Okta и другими identity-провайдерами"),
            ("RBAC scopes (ADR 0021)",
             "granular доступ для пользователей, агентов и токенов"),
            ("Tenant isolation (ADR 0020)",
             "разделение данных, секретов и сетевого периметра по тенантам"),
            ("Подписанные push-уведомления (ADR 0009)",
             "ES256 JWS + envelope-шифрование исходящих событий"),
        ],
        [
            ("Готовность к compliance-ревью",
             "контролируемые точки доступа и аудит-следы"),
            ("Безопасное расширение на всю компанию",
             "тенанты и scopes позволяют разводить отделы и регионы"),
            ("Защита данных за пределами периметра",
             "криптография на исходящих интеграциях"),
            ("Меньше теневого AI",
             "разрешённый централизованный путь дешевле, чем запреты"),
        ],
    )

    # 7. Vendor & model agnosticism
    slide_feature(
        prs, 7, "ТРЕБОВАНИЕ 3 · VENDOR-АГНОСТИЧНОСТЬ",
        "Свобода выбора моделей и провайдеров",
        "Multi-LLM каталог в OpenAI-совместимом формате (ADR 0012). "
        "Модель выбирается под задачу, не под контракт с одним поставщиком.",
        [
            ("Multi-LLM каталог (ADR 0012)",
             "OpenAI, Anthropic, Google, AWS Bedrock, локальные модели"),
            ("OpenAI-совместимый интерфейс",
             "единый формат запросов, прозрачная маршрутизация"),
            ("Native LLM tool-calling (ADR 0011)",
             "tool-calling работает одинаково для всех совместимых провайдеров"),
            ("Гексагональные адаптеры",
             "новый провайдер = новый адаптер, без переписывания ядра"),
        ],
        [
            ("Защита прошлых инвестиций в AI",
             "прежние интеграции и промпты переживают смену провайдера"),
            ("Лучшая модель под каждую задачу",
             "дешёвая для рутины, мощная — для критичных шагов"),
            ("Переговорная позиция с вендорами",
             "возможность мигрировать = аргумент в коммерции"),
            ("Региональные ограничения",
             "выбор провайдера по геозоне и data-residency требованиям"),
        ],
    )

    # 8. Operational efficiency & cost
    slide_feature(
        prs, 8, "ТРЕБОВАНИЕ 4 · СТОИМОСТЬ И ПРОИЗВОДИТЕЛЬНОСТЬ",
        "Контроль токенов и операционных расходов",
        "Токены и compute растут нелинейно. ork снижает расход за счёт "
        "DAG-исполнителя, dynamic embeds и filtered MCP-контекста.",
        [
            ("Workflow engine (ADR 0018)",
             "DAG с параллельным исполнением, без лишних круговых вызовов LLM"),
            ("Dynamic embeds (ADR 0015)",
             "в промпт подставляются только вычисленные значения, не сырые данные"),
            ("Native + LLM-tools + MCP (ADR 0010, 0011)",
             "три слоя инструментов: дешёвые шаги выполняются без LLM"),
            ("Rust runtime",
             "низкий overhead на агента, нет GC-пауз, плотная упаковка процессов"),
        ],
        [
            ("Расход токенов ↓",
             "в LLM попадает только релевантный контекст"),
            ("Время ответа ↓",
             "параллельный DAG быстрее последовательной цепочки"),
            ("Предсказуемая стоимость пилота",
             "видно, какой шаг сколько стоит"),
            ("Меньше железа на ту же нагрузку",
             "Rust позволяет уплотнить инстансы"),
        ],
    )

    # 9. Scalability & resilience
    slide_feature(
        prs, 9, "ТРЕБОВАНИЕ 5 · МАСШТАБ И УСТОЙЧИВОСТЬ",
        "Production-ready event-driven архитектура",
        "Sync-нагрузка через Kong, async — через Kafka. Агенты исполняются "
        "параллельно, отказ одного компонента не валит остальные.",
        [
            ("Гибридный транспорт (ADR 0004)",
             "Kong/HTTP+SSE для sync, Kafka для async — горизонтальное масштабирование"),
            ("Параллельный DAG (ADR 0018)",
             "независимые шаги workflow исполняются одновременно"),
            ("Decoupled адаптеры",
             "сбой gateway, MCP-tool или LLM локализован, не каскадирует"),
            ("Push с retry (ADR 0009)",
             "доставка событий гарантированно или явно отложена"),
        ],
        [
            ("Масштаб под рост нагрузки",
             "добавление инстансов, не переписывание архитектуры"),
            ("SLA по сценариям",
             "graceful degradation при отказе одного провайдера"),
            ("Готовность к пиковым нагрузкам",
             "Kafka буферизует всплески async-задач"),
            ("Production-зрелость, а не demo",
             "архитектура принята как ADR, а не подобрана случайно"),
        ],
    )

    # 10. Observability & governance
    slide_feature(
        prs, 10, "ТРЕБОВАНИЕ 6 · OBSERVABILITY И GOVERNANCE",
        "Единая точка контроля над AI-действиями",
        "Workflow viewer и трассировка (ADR 0022 в работе) дают сквозной обзор: "
        "кто кого вызвал, какой LLM использован, сколько потрачено.",
        [
            ("Workflow viewer (ADR 0017)",
             "real-time просмотр шагов, parts, артефактов и LLM-вызовов"),
            ("Трассировка A2A (ADR 0008)",
             "task-id связывает все промежуточные события в один поток"),
            ("Observability (ADR 0022)",
             "OpenTelemetry-метрики, токены и стоимость по сценариям"),
            ("Структурные логи и аудит",
             "все вызовы агентов, MCP-tools и push доступны для compliance"),
        ],
        [
            ("End-to-end видимость",
             "от запроса пользователя до конечного действия"),
            ("Аудит для регуляторов",
             "кто, когда, какой агент, какая модель, какой результат"),
            ("Быстрая диагностика инцидентов",
             "видно, на каком шаге workflow что пошло не так"),
            ("Контроль расходов",
             "стоимость каждого workflow видна как метрика"),
        ],
    )

    # 11. Channels / accessibility
    slide_feature(
        prs, 11, "ТРЕБОВАНИЕ 7 · КАНАЛЫ ДОСТУПА",
        "AI там, где работает пользователь",
        "Generic Gateways (ADR 0013) дают единый адаптерный слой: один и тот же "
        "агент доступен через Slack, Teams, REST или собственный UI.",
        [
            ("Generic Gateways (ADR 0013)",
             "REST, webhook, WebUI, MCP-as-gateway — за одним абстракционным слоем"),
            ("Web UI (ADR 0017)",
             "React-чат-клиент поверх A2A-стримов из коробки"),
            ("Slack / Teams gateway",
             "адаптеры на тот же общий контракт A2A"),
            ("Public API",
             "JSON-RPC + SSE, документированный, версионируемый"),
        ],
        [
            ("Высокая adoption",
             "пользователи остаются в привычных инструментах"),
            ("Единая логика, разные каналы",
             "новый канал = новый gateway, не новый агент"),
            ("Embed в существующие workflows",
             "AI встраивается в текущие процессы, не ломая их"),
            ("Контролируемые точки входа",
             "каждый gateway можно ограничить scope-ом и rate-limit-ом"),
        ],
    )

    # 12. Rapid experimentation
    slide_feature(
        prs, 12, "ТРЕБОВАНИЕ 8 · СКОРОСТЬ ЭКСПЕРИМЕНТОВ",
        "Pilot → production за недели, а не кварталы",
        "Двухконтурная разработка: YAML-конфиги для бизнес-аналитиков, Rust/Python — "
        "для разработчиков. Federation позволяет переиспользовать существующих агентов.",
        [
            ("YAML workflow templates",
             "сценарии описываются декларативно в workflow-templates/"),
            ("Web UI чат (ADR 0017)",
             "быстрая ручная проверка агентов и workflow"),
            ("Federation (ADR 0007)",
             "Python-агенты на LangGraph подключаются как A2A-пиры"),
            ("ADR-driven процесс",
             "каждое решение зафиксировано — новые члены команды быстро онбордятся"),
        ],
        [
            ("Сокращение цикла идея → пилот",
             "недели вместо кварталов на простой сценарий"),
            ("Закрытие skills-gap",
             "бизнес-аналитики работают в YAML, разработчики — в коде"),
            ("Переиспользование инвестиций",
             "существующие LangGraph / Python-агенты не выбрасываются"),
            ("Прозрачная история решений",
             "ADR — документация, которая не устаревает"),
        ],
    )

    # 13–18. closing
    slide_differentiators(prs)
    slide_maturity(prs)
    slide_roadmap(prs)
    slide_business_summary(prs)
    slide_demo(prs)
    slide_call_to_action(prs)

    out = "ork-for-managers-ru.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides.")


if __name__ == "__main__":
    build()
